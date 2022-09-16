#!/usr/bin/env python

import datetime
import json
import pytz
import re
import requests
from pptx import Presentation


def json_print(obj):
    text = json.dumps(obj, sort_keys=True, indent=4)
    print(text)


prs = Presentation('whats_new.pptx')
section_slide_layout = prs.slide_layouts[4]
content_slide_layout = prs.slide_layouts[8]

begin_date = datetime.datetime(2022, 3, 30, 0, 0, 0, 0, pytz.UTC)
end_date = datetime.datetime.now(datetime.timezone.utc) + datetime.timedelta(days=1)

# important, good to know, less interesting
categories = {'i': [], 'g': [], 'l': []}

params = {
    'item.directoryId': 'whats-new',
    'sort_by': 'item.additionalFields.postDateTime',
    'sort_order': 'desc',
    'size': '400',
    'item.locale': 'en_US',
    'tags.id': 'whats-new#year#2022',
    'page': '0'
}

response = requests.get('https://aws.amazon.com/api/dirs/items/search/', params)
json_response = response.json()

metadata = json_response['metadata']
items = json_response['items']

print(metadata)

item_count = 1
for item in items:
    item_data = item['item']
    print(item_data)
    date_text = item_data['dateUpdated']
    entry_date = datetime.datetime.strptime(date_text, "%Y-%m-%dT%H:%M:%S%z")

    if entry_date < begin_date:
        break
    if entry_date > end_date:
        continue

    item_title = item_data['additionalFields']['headline']

    while True:
        try:
            cat = input(str(item_count) + ': ' + item_title + '\n(i)mportant,(g)ood to know,(l)ess interesting?: ').strip()
            item_count += 1
            if cat not in categories:
                raise Exception("Not valid category")
            break
        except:
            pass

    categories[cat].append(item_data)

for cat in ['i', 'g', 'l']:
    slide = prs.slides.add_slide(section_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    if cat == "i":
        text = "Important"
    elif cat == "g":
        text = "Good to know"
    else:
        text = "Less Relevant"
    title_shape.text = text

    for item in categories[cat]:
        date_text = item['dateUpdated']
        entry_date = datetime.datetime.strptime(date_text, "%Y-%m-%dT%H:%M:%S%z")

        time_string = entry_date.strftime("%d %b")

        title = item['additionalFields']['headline']
        description = time_string + "\n\n" + item['additionalFields']['postSummary']

        TAG_RE = re.compile(r'<[^>]+>')

        description = TAG_RE.sub('', description.strip())

        slide = prs.slides.add_slide(content_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        # title_shape.text = title
        p = title_shape.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        hlink = r.hyperlink
        hlink.address = item['additionalFields']['headlineUrl']
        tf = body_shape.text_frame
        tf.text = description

prs.save("output.pptx")
