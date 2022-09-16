#!/usr/bin/env python

import datetime
import time
from splinter import Browser
import sys

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation('whats_new.pptx')
section_slide_layout = prs.slide_layouts[4]
content_slide_layout = prs.slide_layouts[8]

begindate = datetime.datetime(2021, 9, 19)
enddate = datetime.datetime.now() + datetime.timedelta(days=1)

print(begindate, enddate)

# browser = Browser('chrome')
# browser.visit('https://aws.amazon.com/about-aws/whats-new/2021/')
#
# # important, good to know, less interesting
# categories = {'i': [], 'g': [], 'l': []}
#
# mylist = browser.find_by_css('.directory-item')
# for item in mylist:
#     datetext = item.find_by_css('.date')[0].text
#     entrydate = datetime.datetime.strptime(datetext, "Posted On: %b %d, %Y")
#     print(entrydate)
#     if entrydate < begindate:
#         break
#     if entrydate > enddate:
#         continue
#     title = item.find_by_tag('h3')[0].text
#     while True:
#         try:
#             cat = input(title + "\n(i)mportant,(g)ood to know,(l)ess interesting?: ").rstrip()
#             if cat not in categories:
#                 raise Exception("Not valid category")
#             break
#         except:
#             pass
#
#     categories[cat].append(item)

for cat in ['i', 'g', 'l']:
    slide = prs.slides.add_slide(section_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    if cat == "i":
        text = "Important"
    elif cat == "g":
        text = "Good to know"
    else:
        text = "Less Relevant"
    title_shape.text = text

    # for item in categories[cat]:
    #     datetext = item.find_by_css('.date')[0].text
    #     entrydate = datetime.datetime.strptime(datetext, "Posted On: %b %d, %Y")
    #     timestring = entrydate.strftime("%d %b")
    #     title = item.find_by_tag('h3')[0].text
    #     description = timestring + "\n\n" + item.find_by_css('.description')[0].text
    #     slide = prs.slides.add_slide(content_slide_layout)
    #     shapes = slide.shapes
    #     title_shape = shapes.title
    #     body_shape = shapes.placeholders[1]
    #
    #     # title_shape.text = title
    #     p = title_shape.text_frame.paragraphs[0]
    #     r = p.add_run()
    #     r.text = title
    #     hlink = r.hyperlink
    #     hlink.address = item.find_by_tag('a')[0]["href"]
    #     tf = body_shape.text_frame
    #     tf.text = description
# browser.quit()
prs.save("output.pptx")
